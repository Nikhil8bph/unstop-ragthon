import os
from typing import Optional, List
from langchain_ollama import OllamaEmbeddings
from langchain_chroma import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
import docx
from pptx import Presentation
import pandas as pd
import io
import fitz # PyMuPDF
import base64
import requests
import json

OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
CHROMA_PATH = "./chroma_db"

embeddings = OllamaEmbeddings(
    base_url=OLLAMA_BASE_URL,
    model="qwen3-embedding:4b"
)

# Initialize ChromaDB locally
vectorstore = Chroma(
    collection_name="rag_collection",
    embedding_function=embeddings,
    persist_directory=CHROMA_PATH
)

def ingest_document_chunks(chunks_with_metadata: List[dict]):
    """Ingests chunks where each dict has 'text' and 'metadata'."""
    texts = [c['text'] for c in chunks_with_metadata]
    metadatas = [c['metadata'] for c in chunks_with_metadata]
    
    vectorstore.add_texts(texts=texts, metadatas=metadatas)
    return len(texts)

def extract_text_from_docx(content: bytes) -> str:
    doc = docx.Document(io.BytesIO(content))
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_csv(content: bytes) -> str:
    df = pd.read_csv(io.BytesIO(content))
    return df.to_string(index=False)

def analyze_image(content: bytes) -> str:
    """Uses qwen3-vl via Ollama for both OCR and visual description."""
    encoded_image = base64.b64encode(content).decode('utf-8')
    
    payload = {
        "model": "qwen3-vl:4b",
        "prompt": "Read all the text in this image and describe the visual content in detail. Format your response as a comprehensive description suitable for document retrieval.",
        "stream": False,
        "images": [encoded_image]
    }
    
    try:
        response = requests.post(f"{OLLAMA_BASE_URL}/api/generate", json=payload)
        response.raise_for_status()
        return response.json().get("response", "No analysis available.")
    except Exception as e:
        print(f"Ollama vision error: {e}")
        return "Image analysis failed."

def extract_text_from_pdf(content: bytes) -> List[dict]:
    """Returns list of chunks with text and page metadata."""
    doc = fitz.open(stream=content, filetype="pdf")
    pages_data = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        if text.strip():
            pages_data.append({
                "text": text,
                "page": page_num + 1
            })
    return pages_data

def query_rag(query: str, project_id: Optional[str] = None, folder_id: Optional[str] = None, file_id: Optional[str] = None):
    """
    Search vector DB with metadata filters.
    """
    filter_dict = {}
    if file_id:
        filter_dict["file_id"] = str(file_id)
    elif folder_id:
        filter_dict["folder_id"] = str(folder_id)
    elif project_id:
        filter_dict["project_id"] = str(project_id)
        
    results = vectorstore.similarity_search(query, k=5, filter=filter_dict if filter_dict else None)
    return results
